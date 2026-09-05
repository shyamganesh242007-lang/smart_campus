import os
import io
import re
import time
import math
import asyncio
import difflib
import traceback
from collections import Counter
from dotenv import load_dotenv

# Ensure backend/.env is loaded regardless of working directory
backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
load_dotenv(dotenv_path=os.path.join(backend_dir, ".env"))
load_dotenv()

import docx
import pptx
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_google_genai import ChatGoogleGenerativeAI

try:
    from rapidfuzz import process as rf_process, fuzz as rf_fuzz
    RAPIDFUZZ_AVAILABLE = True
except ImportError:
    RAPIDFUZZ_AVAILABLE = False

try:
    from rank_bm25 import BM25Okapi as RankBM25Okapi
    RANK_BM25_AVAILABLE = True
except ImportError:
    RANK_BM25_AVAILABLE = False


# Comprehensive typo dictionary for automatic query correction
COMMON_TYPOS = {
    "wat": "what",
    "wht": "what",
    "dbmss": "dbms",
    "dbm": "dbms",
    "sqll": "sql",
    "pyhton": "python",
    "pythn": "python",
    "distributd": "distributed",
    "databse": "database",
    "databses": "databases",
    "databas": "database",
    "relatonal": "relational",
    "artifical": "artificial",
    "artificl": "artificial",
    "inteligence": "intelligence",
    "intelligenc": "intelligence",
    "softare": "software",
    "programing": "programming",
    "explian": "explain",
    "explan": "explain",
    "defne": "define",
    "archtecture": "architecture",
    "architecure": "architecture",
    "injecton": "injection",
    "injekshun": "injection",
    "fragmentaton": "fragmentation",
    "fragmntation": "fragmentation",
    "normlization": "normalization",
    "normalizaton": "normalization",
    "securty": "security",
    "vulnerbility": "vulnerability",
    "algoritm": "algorithm",
    "algorthm": "algorithm",
    "decison": "decision",
    "machne": "machine",
    "learnng": "learning",
    "learnin": "learning",
    "photosynthes": "photosynthesis",
    "photosynths": "photosynthesis",
    "biologi": "biology",
    "quantm": "quantum",
    "computng": "computing",
    "looping": "loops",
    "loopss": "loops",
}

# Domain acronyms
DEFAULT_ACRONYMS = {
    "dbms": "database management system",
    "rdbms": "relational database management system",
    "ddbms": "distributed database management system",
    "sql": "structured query language",
    "nosql": "not only sql database",
    "ai": "artificial intelligence",
    "nlp": "natural language processing",
    "rag": "retrieval augmented generation",
    "api": "application programming interface",
    "ml": "machine learning",
    "dl": "deep learning",
    "acid": "atomicity consistency isolation durability",
    "crud": "create read update delete",
    "os": "operating system",
    "gui": "graphical user interface",
    "jvm": "java virtual machine",
    "oop": "object oriented programming",
    "rbac": "role based access control",
    "http": "hypertext transfer protocol",
    "rest": "representational state transfer",
    "json": "javascript object notation",
    "html": "hypertext markup language",
    "css": "cascading style sheets",
    "iot": "internet of things",
    "cpu": "central processing unit",
    "gpu": "graphics processing unit",
    "ram": "random access memory",
    "rom": "read only memory",
}

BASE_VOCABULARY = set(DEFAULT_ACRONYMS.keys()) | set(COMMON_TYPOS.values()) | {
    "what", "is", "explain", "describe", "define", "how", "does", "python", "java", "database",
    "distributed", "injection", "fragmentation", "table", "schema", "query", "index", "system",
    "network", "security", "data", "software", "relational", "architecture", "attribute", "relation",
    "inheritance", "decorators", "polymorphism", "encapsulation", "abstraction", "function", "variable",
    "algorithm", "complexity", "structure", "protocol", "interface", "component", "pipeline", "model",
    "machine", "learning", "photosynthesis", "biology", "plants", "chlorophyll", "loops", "loop",
    "quantum", "computing", "artificial", "intelligence", "transactions", "acid", "isolation"
}

STOP_WORDS = {
    "what", "is", "the", "explain", "about", "how", "does", "are", "tell", "me",
    "a", "an", "in", "of", "and", "or", "to", "for", "with", "on", "at", "by", "from"
}


def extract_marks_and_clean_query(question: str) -> tuple[int | None, str]:
    """
    Extracts examination mark allocations from natural typed queries without any UI selector.
    Supports patterns such as:
      - '2 mark define DBMS'
      - '8 marks explain SQL Injection'
      - '16 mark differentiate DBMS and RDBMS'
      - '10 marks explain Photosynthesis'
      - '5m explain ACID properties'
      - 'what is DBMS (2 marks)'
      - 'explain normalization for 8 marks'
    Returns:
      (marks: int | None, clean_question: str)
    """
    if not question:
        return None, ""

    clean_q = question.strip()
    marks = None

    # 1. Parenthesized marks: e.g. (2 marks), [8 marks], (16m), (5 mark)
    paren_match = re.search(r'[\(\[\{]\s*(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*[\)\]\}]', clean_q, re.IGNORECASE)
    if paren_match:
        marks = int(paren_match.group(1))
        clean_q = re.sub(r'[\(\[\{]\s*(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*[\)\]\}]', ' ', clean_q, flags=re.IGNORECASE)

    # 2. Leading marks: e.g. "2 mark define...", "8 marks explain...", "16m what is..."
    if not marks:
        leading_match = re.search(r'^(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*[-:]?\s+', clean_q, re.IGNORECASE)
        if leading_match:
            marks = int(leading_match.group(1))
            clean_q = re.sub(r'^(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*[-:]?\s+', '', clean_q, flags=re.IGNORECASE)

    # 3. Trailing marks: e.g. "... for 8 marks", "... 16 marks", "... 2m"
    if not marks:
        trailing_match = re.search(r'\s+[-:]?\s*(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*$', clean_q, re.IGNORECASE)
        if trailing_match:
            marks = int(trailing_match.group(1))
            clean_q = re.sub(r'\s+[-:]?\s*(?:for\s*)?(\d{1,2})\s*(?:marks?|mark|m)\s*$', '', clean_q, flags=re.IGNORECASE)

    # 4. In-between marks: e.g. "explain for 8 marks SQL Injection"
    if not marks:
        mid_match = re.search(r'\b(?:for\s+)?(\d{1,2})\s*(?:marks?|mark)\b', clean_q, re.IGNORECASE)
        if mid_match:
            marks = int(mid_match.group(1))
            clean_q = re.sub(r'\b(?:for\s+)?(\d{1,2})\s*(?:marks?|mark)\b', ' ', clean_q, flags=re.IGNORECASE)

    # Clean up whitespace residue
    clean_q = re.sub(r'\s+', ' ', clean_q).strip()
    return marks, clean_q


def detect_question_intent(question: str) -> str:
    """
    Classifies question intent to guide document chunk reranking and response structuring.
    Returns: 'comparison' | 'advantages' | 'disadvantages' | 'types' | 'architecture' | 'process' | 'reasons' | 'definition' | 'general'
    """
    if not question:
        return "general"
    q = question.lower()

    # 1. Comparison / Differences
    if re.search(r'\b(?:differentiate|difference between|differences between|compare|vs|versus|comparison|distinguish|distinction)\b', q):
        return "comparison"

    # 2. Advantages / Benefits
    if re.search(r'\b(?:advantages?|benefits?|merits?|pros|plus points|importance)\b', q):
        return "advantages"

    # 3. Disadvantages / Limitations
    if re.search(r'\b(?:disadvantages?|limitations?|drawbacks?|demerits?|cons|challenges)\b', q):
        return "disadvantages"

    # 4. Types / Classification
    if re.search(r'\b(?:types? of|classification of|categories of|kinds of|different types|variants)\b', q):
        return "types"

    # 5. Architecture / Structure / Components
    if re.search(r'\b(?:architecture|components?|structure|layers?|framework|tier|tiers)\b', q):
        return "architecture"

    # 6. Process / Mechanism / Working / Steps
    if re.search(r'\b(?:how does|how do|how to|working of|working principle|mechanism of|steps? (?:in|of|for)|process of|lifecycle|workflow)\b', q):
        return "process"

    # 7. Reasons / Why / Causes
    if re.search(r'\b(?:why\b|need for|purpose of|reasons? (?:for|behind|why)?|causes? (?:of|for)?|why do we|why is)\b', q):
        return "reasons"

    # 8. Definition
    if re.search(r'\b(?:what is|what are|define|definition of|meaning of|stands for|concept of)\b', q):
        return "definition"

    return "general"


def get_mark_guidelines(marks: int | None) -> str:
    """
    Returns strict length directives for LLM generation based on detected marks.
    """
    if marks is None:
        return "- Provide a complete, comprehensive, and well-structured answer."
    elif marks <= 2:
        return (
            "- Strict Length for 2 Marks: Provide a very short, crisp, direct answer (2 to 3 concise bullet points or 2-3 sentences covering the exact definition/core point). "
            "Do NOT write long paragraphs, redundant intro, or excessive details."
        )
    elif marks <= 5:
        return (
            "- Strict Length for 5 Marks: Provide a concise, focused answer (4 to 5 structured bullet points or 1-2 short paragraphs highlighting the main concepts and key details)."
        )
    elif marks <= 8:
        return (
            "- Strict Length for 8 Marks: Provide a medium-length, well-structured answer with headings, key points, and core explanations."
        )
    elif marks <= 10:
        return (
            "- Strict Length for 10 Marks: Provide a detailed, comprehensive answer with clear headings, in-depth explanations, bullet points, and relevant examples or comparisons."
        )
    else:  # 16 marks or higher
        return (
            f"- Strict Length for {marks} Marks: Provide a full descriptive university exam answer with clear introduction, major structured sections/headings, detailed explanations, bullet points, real-world examples, and tables/diagram descriptions where applicable."
        )


def get_intent_guidelines(intent: str) -> str:
    """
    Returns structural directives matching the detected question intent.
    """
    if intent == "comparison":
        return "- Question Intent is COMPARISON: Present the differences clearly using a clean Markdown comparison table or distinct comparative criteria."
    elif intent == "definition":
        return "- Question Intent is DEFINITION: State the formal definition clearly and crisply at the beginning, followed by key characteristics."
    elif intent == "reasons":
        return "- Question Intent is REASONS / JUSTIFICATION: Focus directly on the reasons, motivation, problems solved, and necessity."
    elif intent == "process":
        return "- Question Intent is PROCESS / MECHANISM: Present the steps, workflow, or mechanism in clear, chronological numbered sequence."
    elif intent == "advantages":
        return "- Question Intent is ADVANTAGES: Focus on the key benefits, merits, and positive impacts with clear bullet points."
    elif intent == "disadvantages":
        return "- Question Intent is DISADVANTAGES: Focus on the limitations, drawbacks, and challenges with clear bullet points."
    elif intent == "types":
        return "- Question Intent is TYPES / CLASSIFICATION: Enumerate and explain each type or category with dedicated headings or sub-bullets."
    elif intent == "architecture":
        return "- Question Intent is ARCHITECTURE: Detail the architectural components, tiers, layers, and how data/control flows between them."
    return "- Structure your response with clear headings and bullet points where appropriate."


def merge_two_chunks(chunk_a: str, chunk_b: str) -> str:
    """
    Merges two adjacent document chunks cleanly by removing overlap redundancy.
    Preserves headings, bullet points, numbers, and paragraphs without word duplication.
    """
    if not chunk_a:
        return chunk_b or ""
    if not chunk_b:
        return chunk_a or ""

    max_len = min(len(chunk_a), len(chunk_b), 250)
    best_overlap = 0

    for l in range(max_len, 15, -1):
        if chunk_a[-l:] == chunk_b[:l]:
            best_overlap = l
            break
        if chunk_a[-l:].strip() == chunk_b[:l].strip():
            best_overlap = l
            break

    if best_overlap > 0:
        return chunk_a + chunk_b[best_overlap:]
    
    # Clean separator preserving paragraphs
    if chunk_a.endswith("\n") or chunk_b.startswith("\n"):
        return chunk_a + "\n" + chunk_b.lstrip("\n")
    return chunk_a + "\n\n" + chunk_b


def format_and_split_long_answer(text: str, max_chars: int = 3800) -> str:
    """
    Allows responses up to around 3000-5000 characters before splitting.
    If the answer exceeds max_chars, automatically continues with Part 2, Part 3, etc.
    Never cuts the answer in the middle of a sentence.
    Preserves headings, bullet points, numbering, and paragraphs.
    """
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text

    parts = []
    remaining = text

    while len(remaining) > max_chars:
        cut_point = -1
        search_window = remaining[:max_chars]

        # Prioritize natural paragraph breaks, then sentence endings
        patterns = [
            r'\n\n',
            r'\.\s+',
            r'\?\s+',
            r'!\s+',
            r'\n'
        ]
        for pat in patterns:
            matches = list(re.finditer(pat, search_window))
            if matches:
                for m in reversed(matches):
                    if m.end() >= int(max_chars * 0.65):
                        cut_point = m.end()
                        break
            if cut_point != -1:
                break

        # If no clean break found in the lower window, look forward up to 600 chars for sentence end
        if cut_point == -1:
            extended = remaining[:max_chars + 600]
            for pat in [r'\n\n', r'\.\s+', r'\?\s+', r'!\s+']:
                m = re.search(pat, extended[max_chars:])
                if m:
                    cut_point = max_chars + m.end()
                    break

        if cut_point == -1:
            cut_point = max_chars

        part_text = remaining[:cut_point].strip()
        parts.append(part_text)
        remaining = remaining[cut_point:].strip()

    if remaining:
        parts.append(remaining)

    if len(parts) <= 1:
        return text

    formatted_parts = []
    for i, p in enumerate(parts, 1):
        formatted_parts.append(f"### Part {i}\n\n{p}")

    return "\n\n---\n\n".join(formatted_parts)


class FastBM25:
    """In-memory BM25Okapi implementation for sparse keyword retrieval."""
    def __init__(self, tokenized_corpus, k1=1.5, b=0.75):
        self.k1 = k1
        self.b = b
        self.corpus_size = len(tokenized_corpus)
        self.doc_lens = [len(doc) for doc in tokenized_corpus]
        self.avgdl = sum(self.doc_lens) / self.corpus_size if self.corpus_size > 0 else 0
        self.doc_freqs = []
        self.nd = Counter()

        for doc in tokenized_corpus:
            freq = Counter(doc)
            self.doc_freqs.append(freq)
            for word in freq.keys():
                self.nd[word] += 1

        self.idf = {}
        for word, freq in self.nd.items():
            self.idf[word] = math.log((self.corpus_size - freq + 0.5) / (freq + 0.5) + 1.0)

    def get_scores(self, query_tokens):
        scores = [0.0] * self.corpus_size
        if self.corpus_size == 0 or self.avgdl == 0:
            return scores

        for q in query_tokens:
            if q not in self.idf:
                continue
            idf_val = self.idf[q]
            for i, doc_freq in enumerate(self.doc_freqs):
                if q in doc_freq:
                    freq = doc_freq[q]
                    doc_len = self.doc_lens[i]
                    numerator = idf_val * freq * (self.k1 + 1)
                    denominator = freq + self.k1 * (1 - self.b + self.b * (doc_len / self.avgdl))
                    scores[i] += (numerator / denominator)
        return scores


class CachedEmbeddings:
    """Embeddings wrapper with query and document caching to prevent redundant calculations."""
    def __init__(self, base_embeddings):
        self.base = base_embeddings
        self.cache = {}

    def embed_documents(self, texts):
        uncached = [t for t in texts if t not in self.cache]
        if uncached:
            embeddings = self.base.embed_documents(uncached)
            for t, emb in zip(uncached, embeddings):
                self.cache[t] = emb
        return [self.cache[t] for t in texts]

    def embed_query(self, text):
        clean_key = text.strip().lower()
        if clean_key not in self.cache:
            self.cache[clean_key] = self.base.embed_query(text)
        return self.cache[clean_key]


class RAGService:
    def __init__(self):
        self.api_key = os.getenv("GEMINI_API_KEY")
        self.ocr_reader = None
        self.persist_dir = "./.chroma_db"
        self.collection_name = "studygen_active_session"
        self.uploaded_files = {}

        # Automatically ensure required directories exist
        os.makedirs(self.persist_dir, exist_ok=True)
        os.makedirs("uploads", exist_ok=True)
        os.makedirs("backend/uploads", exist_ok=True)
        os.makedirs("backend/.chroma_db", exist_ok=True)

        # In-memory indices & caches
        self.corpus_chunks = []
        self.dynamic_vocabulary = set()
        self.dynamic_acronyms = {}
        self.bm25_index = None
        self.response_cache = {}

        try:
            base_embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2"
            )
            self.embeddings = CachedEmbeddings(base_embeddings)

            import chromadb
            self.chroma_client = chromadb.PersistentClient(path=self.persist_dir)
            self.chroma_client.get_or_create_collection(
                self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
            self.vectorstore = Chroma(
                client=self.chroma_client,
                collection_name=self.collection_name,
                embedding_function=self.embeddings,
                collection_metadata={"hnsw:space": "cosine"}
            )

            # Load active session documents on startup
            self._reload_corpus_from_store()

        except Exception as e:
            print("Failed to initialize embeddings or ChromaDB:")
            traceback.print_exc()
            self.embeddings = None
            self.vectorstore = None

        # Semantic chunking: 700 size, 150 overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=700,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", "? ", "! ", "; ", " ", ""]
        )

    def _get_ocr_reader(self):
        if self.ocr_reader is None:
            import easyocr
            self.ocr_reader = easyocr.Reader(['en'], gpu=False)
        return self.ocr_reader

    def _reload_corpus_from_store(self):
        """Loads stored chunks from ChromaDB into memory and rebuilds BM25, vocab, and acronyms."""
        self.corpus_chunks = []
        self.dynamic_vocabulary = set()
        self.dynamic_acronyms = {}
        self.uploaded_files = {}
        self.response_cache.clear()

        try:
            if not self.vectorstore:
                return
            count = self.vectorstore._collection.count()
            if count > 0:
                results = self.vectorstore._collection.get(include=["documents", "metadatas"])
                docs = results.get("documents", [])
                metas = results.get("metadatas", [])

                for doc_text, meta in zip(docs, metas):
                    if not doc_text:
                        continue
                    m = meta or {}
                    fn = m.get("filename", "unknown")
                    ft = m.get("file_type", m.get("fileType", ""))

                    if fn not in self.uploaded_files:
                        self.uploaded_files[fn] = {
                            "fileType": ft,
                            "total_chunks": 0,
                            "chunkCount": 0,
                            "timestamp": float(m.get("upload_timestamp", time.time()))
                        }
                    self.uploaded_files[fn]["total_chunks"] += 1
                    self.uploaded_files[fn]["chunkCount"] += 1

                    self.corpus_chunks.append({
                        "text": doc_text,
                        "metadata": m
                    })

                    # Populate dynamic vocabulary
                    words = re.findall(r'[a-zA-Z]{3,}', doc_text.lower())
                    self.dynamic_vocabulary.update(words)

                    # Extract dynamic acronyms
                    doc_acronyms = self._extract_acronyms(doc_text)
                    self.dynamic_acronyms.update(doc_acronyms)

                self._rebuild_bm25()
                print(f"[StudyGen RAG] Loaded {len(self.uploaded_files)} active files ({len(self.corpus_chunks)} chunks) into Hybrid Index.")
        except Exception as e:
            print(f"[StudyGen RAG] Notice while reloading corpus from store: {e}")

    def _extract_acronyms(self, text: str) -> dict:
        """Extracts acronym definitions dynamically from uploaded document text."""
        acronyms = {}
        try:
            # Pattern 1: Full Name (ACRONYM) e.g., "Relational Database Management System (RDBMS)"
            p1 = re.compile(r'\b([A-Z][a-zA-Z0-9]+(?:\s+[A-Z][a-zA-Z0-9]+|\s+(?:of|and|for|in|the)\s+[A-Z][a-zA-Z0-9]+){1,5})\s*\(([A-Z]{2,8})\)')
            for full_name, acr in p1.findall(text):
                acr_clean = acr.strip().lower()
                full_clean = full_name.strip().lower()
                acronyms[acr_clean] = full_clean

            # Pattern 2: ACRONYM (Full Name) e.g., "DBMS (Database Management System)"
            p2 = re.compile(r'\b([A-Z]{2,8})\s*\(([A-Za-z\s]{4,60})\)')
            for acr, full_name in p2.findall(text):
                acr_clean = acr.strip().lower()
                full_clean = full_name.strip().lower()
                if len(full_clean.split()) >= 2:
                    acronyms[acr_clean] = full_clean

            # Pattern 3: ACRONYM: Full Name or ACRONYM - Full Name
            p3 = re.compile(r'\b([A-Z]{2,8})\s*[:\-–]\s*([A-Z][a-zA-Z\s]{4,50})')
            for acr, full_name in p3.findall(text):
                acr_clean = acr.strip().lower()
                full_clean = full_name.strip().lower()
                if len(full_clean.split()) >= 2:
                    acronyms[acr_clean] = full_clean
        except Exception as e:
            print(f"[StudyGen RAG] Error extracting acronyms: {e}")

        return acronyms

    def _rebuild_bm25(self):
        """Rebuilds the BM25 index over the current active corpus chunks."""
        if not self.corpus_chunks:
            self.bm25_index = None
            return

        tokenized_corpus = [
            [w for w in re.findall(r'\w+', item["text"].lower()) if len(w) > 1 and w not in STOP_WORDS]
            for item in self.corpus_chunks
        ]

        if RANK_BM25_AVAILABLE:
            try:
                self.bm25_index = RankBM25Okapi(tokenized_corpus)
                return
            except Exception:
                pass

        self.bm25_index = FastBM25(tokenized_corpus)

    def normalize_question(self, question: str) -> str:
        """
        Normalizes query using RapidFuzz typo correction against base & dynamic vocabulary,
        and automatically expands acronyms from uploaded documents.
        """
        q = question.lower().strip()
        q_clean = re.sub(r'[^\w\s]', ' ', q)
        raw_words = q_clean.split()

        all_vocabulary = BASE_VOCABULARY | self.dynamic_vocabulary
        all_acronyms = {**DEFAULT_ACRONYMS, **self.dynamic_acronyms}

        corrected_words = []
        for word in raw_words:
            if word in COMMON_TYPOS:
                corrected_words.append(COMMON_TYPOS[word])
            elif word in all_acronyms or word in all_vocabulary:
                corrected_words.append(word)
            elif len(word) >= 3:
                # RapidFuzz typo correction
                corrected = word
                if RAPIDFUZZ_AVAILABLE and all_vocabulary:
                    match = rf_process.extractOne(
                        word,
                        all_vocabulary,
                        scorer=rf_fuzz.ratio,
                        score_cutoff=75
                    )
                    if match:
                        corrected = match[0]
                else:
                    matches = difflib.get_close_matches(word, all_vocabulary, n=1, cutoff=0.75)
                    if matches:
                        corrected = matches[0]
                corrected_words.append(corrected)
            else:
                corrected_words.append(word)

        # Automatic acronym expansion
        expanded_words = []
        for word in corrected_words:
            expanded_words.append(word)
            if word in all_acronyms:
                expanded_words.append(all_acronyms[word])

        return " ".join(expanded_words).strip()

    def process_file(self, file_bytes: bytes, filename: str) -> str:
        """
        Universal file support: PDF, DOCX, PPTX, TXT, MD.
        Applies semantic chunking (700 size, 150 overlap) and updates dynamic vocabulary,
        acronyms, ChromaDB, and BM25 index. Prevents duplicate embeddings.
        """
        if not self.embeddings or not self.vectorstore:
            raise Exception("Vectorstore service is not initialized.")

        try:
            ext = filename.split('.')[-1].lower() if '.' in filename else ''
            structured_chunks = []

            # 1. Universal file extraction
            if ext == 'pdf':
                reader = PdfReader(io.BytesIO(file_bytes))
                for page_idx, page in enumerate(reader.pages):
                    extracted = page.extract_text()
                    if extracted and extracted.strip():
                        page_chunks = self.text_splitter.split_text(extracted.strip())
                        for chunk in page_chunks:
                            structured_chunks.append({"text": chunk, "page": page_idx + 1})

            elif ext == 'docx':
                doc = docx.Document(io.BytesIO(file_bytes))
                full_text = ""
                for para in doc.paragraphs:
                    p_text = para.text.strip()
                    if p_text:
                        style_name = para.style.name.lower() if para.style else ""
                        if "heading" in style_name:
                            full_text += f"\n## {p_text}\n\n"
                        elif any(kw in style_name for kw in ["list", "bullet", "number"]):
                            full_text += f"- {p_text}\n"
                        else:
                            full_text += f"{p_text}\n\n"

                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        dedup = []
                        for c in row_cells:
                            if not dedup or c != dedup[-1]:
                                dedup.append(c)
                        if dedup:
                            full_text += " | ".join(dedup) + "\n"
                    full_text += "\n"

                chunks = self.text_splitter.split_text(full_text)
                for chunk in chunks:
                    structured_chunks.append({"text": chunk})

            elif ext in ['pptx', 'ppt']:
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                for slide_idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    if slide_texts:
                        slide_content = f"Slide {slide_idx + 1}:\n" + "\n".join(slide_texts)
                        slide_chunks = self.text_splitter.split_text(slide_content)
                        for chunk in slide_chunks:
                            structured_chunks.append({"text": chunk, "slide": slide_idx + 1})

            elif ext in ['txt', 'md', 'markdown']:
                try:
                    raw_text = file_bytes.decode('utf-8')
                except UnicodeDecodeError:
                    raw_text = file_bytes.decode('latin1', errors='ignore')
                chunks = self.text_splitter.split_text(raw_text)
                for chunk in chunks:
                    structured_chunks.append({"text": chunk})

            elif ext in ['png', 'jpg', 'jpeg', 'webp']:
                reader = self._get_ocr_reader()
                result = reader.readtext(file_bytes, detail=0)
                raw_text = " ".join(result)
                chunks = self.text_splitter.split_text(raw_text)
                for chunk in chunks:
                    structured_chunks.append({"text": chunk})

            else:
                raise Exception(f"Unsupported file format: {ext}")

            if not structured_chunks:
                raise Exception("Unable to extract readable content from this file.")

            doc_id = f"{filename}_{int(time.time())}"
            texts_to_add = [item["text"] for item in structured_chunks]

            metadatas = [
                {
                    "filename": filename,
                    "file_type": ext,
                    "document_id": doc_id,
                    "upload_timestamp": str(time.time()),
                    "chunk_index": i,
                    "total_chunks": len(structured_chunks),
                    **({"slide_number": structured_chunks[i]["slide"]} if "slide" in structured_chunks[i] else {}),
                    **({"page_number": structured_chunks[i]["page"]} if "page" in structured_chunks[i] else {})
                }
                for i in range(len(structured_chunks))
            ]

            # Prevent duplicate embeddings: clear previous chunks for this filename
            try:
                self.vectorstore.delete(where={"filename": filename})
            except Exception:
                pass

            # Add to ChromaDB in batches
            batch_size = 50
            for i in range(0, len(texts_to_add), batch_size):
                self.vectorstore.add_texts(
                    texts=texts_to_add[i:i+batch_size],
                    metadatas=metadatas[i:i+batch_size]
                )

            # Update in-memory corpus chunks
            self.corpus_chunks = [c for c in self.corpus_chunks if c.get("metadata", {}).get("filename") != filename]
            for text_chunk, meta in zip(texts_to_add, metadatas):
                self.corpus_chunks.append({"text": text_chunk, "metadata": meta})

                # Dynamic vocabulary expansion
                words = re.findall(r'[a-zA-Z]{3,}', text_chunk.lower())
                self.dynamic_vocabulary.update(words)

                # Automatic acronym extraction from document
                extracted_acrs = self._extract_acronyms(text_chunk)
                self.dynamic_acronyms.update(extracted_acrs)

            # Invalidate query cache when new files are processed
            self.response_cache.clear()

            # Rebuild BM25 index with new chunks
            self._rebuild_bm25()

            self.uploaded_files[filename] = {
                "fileType": ext,
                "chunkCount": len(texts_to_add),
                "total_chunks": len(texts_to_add),
                "timestamp": time.time()
            }

            print(f"[StudyGen RAG] Indexed '{filename}' ({ext.upper()}): {len(texts_to_add)} chunks stored with Semantic Chunking.")
            return ext

        except Exception as e:
            print(f"[StudyGen RAG] Error processing file {filename}: {e}")
            traceback.print_exc()
            raise Exception("Unable to process this file. Please try again.")

    def delete_file(self, filename: str):
        """Deletes file chunks from ChromaDB and in-memory corpus, updating BM25 and vocabulary."""
        if self.vectorstore:
            try:
                self.vectorstore.delete(where={"filename": filename})
                print(f"[StudyGen RAG] Removed chunks for '{filename}'.")
            except Exception as e:
                print(f"[StudyGen RAG] Error removing chunks for {filename}: {e}")

        if filename in self.uploaded_files:
            del self.uploaded_files[filename]

        self.corpus_chunks = [c for c in self.corpus_chunks if c.get("metadata", {}).get("filename") != filename]
        self.response_cache.clear()
        self._rebuild_bm25()

    def _merge_consecutive_document_chunks(self, selected_doc: str, top_results: list) -> str:
        """
        Gathers matching chunks for selected_doc from in-memory corpus,
        expands around relevant chunk indices to retrieve complete contiguous sections,
        and merges consecutive chunks while removing overlap redundancy.
        Preserves headings, bullet points, numbers, and paragraphs.
        """
        doc_chunks = [c for c in self.corpus_chunks if c.get("metadata", {}).get("filename") == selected_doc]
        doc_chunks.sort(key=lambda c: c.get("metadata", {}).get("chunk_index", 0))
        total_doc_chunks = len(doc_chunks)

        if total_doc_chunks == 0:
            return ""

        # Collect relevant chunk indices from top reranked results for this document
        matched_indices = set()
        for res in top_results:
            if res.get("doc_name") == selected_doc:
                c_idx = res.get("metadata", {}).get("chunk_index")
                if c_idx is not None and (res.get("rerank_score", 0) >= 0.20 or res.get("kw_match")):
                    matched_indices.add(c_idx)

        # If none met threshold, take the top candidate for this document
        if not matched_indices:
            for res in top_results:
                if res.get("doc_name") == selected_doc:
                    c_idx = res.get("metadata", {}).get("chunk_index")
                    if c_idx is not None:
                        matched_indices.add(c_idx)
                        break

        if not matched_indices:
            matched_indices.add(0)

        # Expand relevant indices to adjacent chunks to capture full sections/explanations
        expanded_indices = set()
        for idx in matched_indices:
            for offset in [-1, 0, 1, 2]:
                cand = idx + offset
                if 0 <= cand < total_doc_chunks:
                    expanded_indices.add(cand)

        sorted_indices = sorted(expanded_indices)
        if not sorted_indices:
            return doc_chunks[0]["text"]

        # Group contiguous chunk indices into sequential clusters
        clusters = []
        current_cluster = [sorted_indices[0]]
        for idx in sorted_indices[1:]:
            if idx == current_cluster[-1] + 1:
                current_cluster.append(idx)
            else:
                clusters.append(current_cluster)
                current_cluster = [idx]
        clusters.append(current_cluster)

        # Merge contiguous chunks inside each cluster seamlessly
        merged_sections = []
        for cl in clusters:
            cluster_text = doc_chunks[cl[0]]["text"]
            for next_idx in cl[1:]:
                cluster_text = merge_two_chunks(cluster_text, doc_chunks[next_idx]["text"])
            merged_sections.append(cluster_text.strip())

        return "\n\n---\n\n".join(merged_sections).strip()

    async def _call_gemini(self, prompt: str) -> str:
        """
        Executes Gemini generation across active models with zero placeholder responses.
        Uses google.genai Client with fallback to LangChain ChatGoogleGenerativeAI.
        """
        if not self.api_key:
            return ""

        candidate_models = [
            "gemini-3.5-flash-lite",
            "gemini-3.1-flash-lite-preview",
            "gemini-flash-lite-latest"
        ]

        # Fast path: google.genai Client in thread
        try:
            from google import genai
            client = genai.Client(api_key=self.api_key)
            for model_name in candidate_models:
                try:
                    def _sync_gen(m=model_name):
                        res = client.models.generate_content(model=m, contents=prompt)
                        return getattr(res, "text", "") or ""

                    text = await asyncio.wait_for(asyncio.to_thread(_sync_gen), timeout=12.0)
                    if text and text.strip():
                        return text.strip()
                except Exception as m_err:
                    print(f"[StudyGen RAG] genai model '{model_name}' notice: {m_err}")
                    continue
        except Exception as client_err:
            print(f"[StudyGen RAG] google.genai client initialization notice: {client_err}")

        # Secondary path: ChatGoogleGenerativeAI
        for model_name in candidate_models:
            try:
                llm = ChatGoogleGenerativeAI(
                    model=model_name,
                    google_api_key=self.api_key
                )
                resp = await asyncio.wait_for(llm.ainvoke(prompt), timeout=12.0)
                if isinstance(resp.content, list):
                    text = "".join([
                        p if isinstance(p, str) else (p.get("text", "") if isinstance(p, dict) else getattr(p, "text", str(p)))
                        for p in resp.content
                    ]).strip()
                else:
                    text = str(resp.content).strip()

                if text:
                    return text
            except Exception as e:
                print(f"[StudyGen RAG] LangChain model '{model_name}' notice: {e}")
                continue

        return ""

    async def ask_question(self, question: str) -> dict:
        """
        Production Hybrid Retrieval (BM25 + ChromaDB) with Question Intent & Natural Mark Scaling.
        - Naturally detects question marks (e.g. 2 marks, 16 marks, 5m) without dropdowns.
        - Automatically classifies question intent (definition, comparison, process, reasons, etc.)
          and reranks document sections to match.
        - Dual Answer Cards:
          - When document is found: returns both [File Based] (grounded in document) and
            [AI Enhanced] (enriching document with real-world examples and clearer explanation).
          - When document is not found (AI Fallback): returns single [AI Enhanced] card from Gemini AI.
        - Caches repeated queries for instant responses.
        - Guaranteed never to crash or display placeholder errors.
        """
        # Step 0: Extract marks and clean query
        raw_q = question.strip()
        detected_marks, clean_q = extract_marks_and_clean_query(raw_q)
        detected_intent = detect_question_intent(clean_q)
        mark_guidelines = get_mark_guidelines(detected_marks)
        intent_guidelines = get_intent_guidelines(detected_intent)

        print(f"\n[StudyGen RAG] Input: '{raw_q}' | Marks: {detected_marks} | Intent: {detected_intent.upper()} | Clean Q: '{clean_q}'")

        # Check response cache
        active_files_key = tuple(sorted(self.uploaded_files.keys()))
        cache_key = (clean_q.lower(), detected_marks, active_files_key)
        if cache_key in self.response_cache:
            print(f"[StudyGen RAG] Serving cached response for: '{clean_q}' ({detected_marks} marks)")
            return self.response_cache[cache_key]

        # Step 1: Normalize query with RapidFuzz and Acronym Expansion
        normalized_q = self.normalize_question(clean_q)
        print(f"[StudyGen RAG] normalized query: '{normalized_q}'")

        clean_q_terms = [
            t for t in re.sub(r'[^\w\s]', ' ', clean_q.lower()).split()
            if len(t) > 1 and t not in STOP_WORDS
        ]
        norm_q_terms = [
            t for t in re.sub(r'[^\w\s]', ' ', normalized_q.lower()).split()
            if len(t) > 1 and t not in STOP_WORDS
        ]
        all_query_terms = list(dict.fromkeys(clean_q_terms + norm_q_terms))

        candidates = {}

        # Step 2: Dense Retrieval from ChromaDB
        try:
            if self.vectorstore and self.uploaded_files:
                chroma_results = self.vectorstore.similarity_search_with_score(normalized_q, k=8)
                for doc, dist in chroma_results:
                    text = doc.page_content
                    meta = doc.metadata or {}
                    doc_key = (meta.get("filename", ""), meta.get("chunk_index", -1), text[:60])

                    if dist > 1.0:
                        dense_sim = max(0.0, min(1.0, 1.0 - (dist / 2.0)))
                    else:
                        dense_sim = max(0.0, min(1.0, 1.0 - dist))

                    candidates[doc_key] = {
                        "text": text,
                        "metadata": meta,
                        "dense_score": dense_sim,
                        "bm25_score": 0.0
                    }
        except Exception as e:
            print(f"[StudyGen RAG] ChromaDB retrieval error: {e}")

        # Step 3: Sparse Retrieval from BM25 (Hybrid retrieval)
        try:
            if self.bm25_index and self.corpus_chunks and all_query_terms:
                raw_bm25_scores = self.bm25_index.get_scores(all_query_terms)
                max_bm25 = max(raw_bm25_scores) if raw_bm25_scores else 0.0

                top_bm25_indices = sorted(range(len(raw_bm25_scores)), key=lambda i: raw_bm25_scores[i], reverse=True)[:8]
                for idx in top_bm25_indices:
                    raw_sc = raw_bm25_scores[idx]
                    if raw_sc <= 0:
                        continue
                    norm_bm25 = raw_sc / max_bm25 if max_bm25 > 0 else 0.0
                    chunk_item = self.corpus_chunks[idx]
                    text = chunk_item["text"]
                    meta = chunk_item["metadata"] or {}
                    doc_key = (meta.get("filename", ""), meta.get("chunk_index", -1), text[:60])

                    if doc_key in candidates:
                        candidates[doc_key]["bm25_score"] = norm_bm25
                    else:
                        candidates[doc_key] = {
                            "text": text,
                            "metadata": meta,
                            "dense_score": 0.0,
                            "bm25_score": norm_bm25
                        }
        except Exception as e:
            print(f"[StudyGen RAG] BM25 retrieval error: {e}")

        # Step 4: Top-8 Retrieval + Intent-Guided Reranking
        scored_results = []
        for doc_key, data in candidates.items():
            text = data["text"]
            meta = data["metadata"]
            dense_s = data["dense_score"]
            bm25_s = data["bm25_score"]

            doc_lower = text.lower()
            # Exact keyword match
            kw_count = sum(1 for term in all_query_terms if re.search(r'\b' + re.escape(term) + r'\b', doc_lower))
            kw_match = kw_count > 0

            # Exact phrase match bonus
            exact_phrase_match = 1 if (len(clean_q_terms) >= 2 and " ".join(clean_q_terms) in doc_lower) else 0
            kw_boost = min(0.25, 0.08 * kw_count) + (0.15 if exact_phrase_match else 0.0)

            # Intent-guided section matching bonus
            intent_bonus = 0.0
            if detected_intent == "comparison" and re.search(r'\b(?:difference|comparison|vs|differentiate|table|distinction)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "advantages" and re.search(r'\b(?:advantage|benefit|merit|pros)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "disadvantages" and re.search(r'\b(?:disadvantage|limitation|drawback|cons)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "types" and re.search(r'\b(?:type|classification|category|categories|kinds)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "architecture" and re.search(r'\b(?:architecture|component|layer|tier|module|structure)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "process" and re.search(r'\b(?:step|process|working|phase|stage|execution|procedure)\b', doc_lower):
                intent_bonus = 0.15
            elif detected_intent == "definition" and re.search(r'\b(?:defined as|refers to|is a|means)\b', doc_lower):
                intent_bonus = 0.12
            elif detected_intent == "reasons" and re.search(r'\b(?:need|purpose|reason|why|because|in order to)\b', doc_lower):
                intent_bonus = 0.12

            # Combined hybrid reranking score with intent boost
            hybrid_score = (0.50 * dense_s) + (0.35 * bm25_s) + kw_boost + intent_bonus
            rerank_score = min(1.0, hybrid_score)

            scored_results.append({
                "text": text,
                "metadata": meta,
                "doc_name": meta.get("filename", "Unknown"),
                "dense_score": dense_s,
                "bm25_score": bm25_s,
                "kw_match": kw_match,
                "kw_count": kw_count,
                "intent_bonus": intent_bonus,
                "rerank_score": rerank_score
            })

        # Sort candidates descending by rerank score and retain Top-8
        scored_results.sort(key=lambda x: x["rerank_score"], reverse=True)
        top_8_results = scored_results[:8]

        # Step 5: Accurate Document Confidence Evaluation
        top_candidate = top_8_results[0] if top_8_results else None
        top_dense = top_candidate["dense_score"] if top_candidate else 0.0
        top_bm25 = top_candidate["bm25_score"] if top_candidate else 0.0
        top_has_kw = top_candidate["kw_match"] if top_candidate else False
        top_score = top_candidate["rerank_score"] if top_candidate else 0.0

        has_keyword_relevance = (top_has_kw and (top_dense >= 0.20 or top_bm25 > 0.05)) or (top_bm25 >= 0.20)
        has_pure_semantic_high = (top_dense >= 0.68 and top_score >= 0.50)
        is_high_confidence = bool(top_8_results and (
            (has_keyword_relevance and top_score >= 0.28) or has_pure_semantic_high
        ))

        chosen_source = "document" if (is_high_confidence and self.uploaded_files) else "ai"
        selected_source_doc = "Gemini AI"
        matched_file_type = None

        if chosen_source == "document":
            doc_relevance = {}
            for res in top_8_results:
                fn = res["doc_name"]
                doc_relevance[fn] = doc_relevance.get(fn, 0.0) + res["rerank_score"]

            selected_source_doc = max(doc_relevance.items(), key=lambda x: x[1])[0]
            top_meta = next((r["metadata"] for r in top_8_results if r["doc_name"] == selected_source_doc), {})
            matched_file_type = top_meta.get("file_type", "")

        print(f"[StudyGen RAG] Source Decision: {chosen_source.upper()} (Top Score: {top_score:.3f}, KW: {top_has_kw})")

        # Step 6: Dual Answer Generation (File Based & AI Enhanced)
        file_raw_answer = ""
        ai_raw_answer = ""

        if chosen_source == "document":
            merged_doc_text = self._merge_consecutive_document_chunks(selected_source_doc, top_8_results)

            doc_prompt = (
                "You are CampusX StudyGen, an academic study assistant. "
                "Provide an accurate, well-structured answer to the user's question using ONLY the provided document material.\n"
                "Rules:\n"
                "- Base your answer strictly on the document text below.\n"
                f"{mark_guidelines}\n"
                f"{intent_guidelines}\n"
                "- Preserve headings, bullet points, numbering, and structured paragraphs.\n"
                "- Render mathematical equations clearly using LaTeX ($...$ for inline, $$...$$ for display formulas).\n"
                "- Write chemical equations and formulas with clean subscripts and reaction arrows (e.g., CO₂ + H₂O → C₆H₁₂O₆ + O₂).\n"
                "- When presenting structured, comparative, or tabulated information, format it as clean Markdown tables.\n"
                "- Do not include meta statements like 'Based on the document' or 'Source: File Based'.\n"
                "- If the provided text does NOT contain enough information to answer the question, reply EXACTLY with: 'Answer not available in uploaded documents.'\n\n"
                f"Document Content:\n{merged_doc_text}\n\n"
                f"Question:\n{clean_q}"
            )

            gemini_doc_answer = await self._call_gemini(doc_prompt)

            if gemini_doc_answer and "answer not available in uploaded documents" in gemini_doc_answer.lower():
                print("[StudyGen RAG] Document context was insufficient. Falling back immediately to Gemini AI educational answer.")
                chosen_source = "ai"
                selected_source_doc = "Gemini AI"
                matched_file_type = None
            elif gemini_doc_answer and len(gemini_doc_answer.strip()) > 20:
                file_raw_answer = gemini_doc_answer
            else:
                file_raw_answer = merged_doc_text

            # If document context was valid, generate AI Enhanced answer enriched with Gemini knowledge
            if chosen_source == "document":
                ai_enhanced_prompt = (
                    "You are CampusX StudyGen, an expert academic tutor. "
                    "The user asked an academic question, and we retrieved reference material from their course document.\n"
                    "Your task is to generate an 'AI Enhanced' explanation that enriches, clarifies, and elevates the document material.\n"
                    "Rules:\n"
                    "- The document is the primary source of truth; do not contradict it.\n"
                    "- Intelligently combine the document's core concepts with your broader educational explanation: add clear real-world examples, intuitive analogies, missing definitions/context, and improved organization.\n"
                    f"{mark_guidelines}\n"
                    f"{intent_guidelines}\n"
                    "- Render mathematical equations using LaTeX ($...$ for inline, $$...$$ for display formulas).\n"
                    "- Render chemical formulas with clean subscripts (CO₂, H₂O, C₆H₁₂O₆) and reaction arrows (→).\n"
                    "- When presenting comparisons or properties, format them using clean Markdown tables.\n"
                    "- Never include meta statements like 'According to the document' or 'Here is an enhanced explanation'. Directly provide the enhanced academic explanation.\n\n"
                    f"Document Reference:\n{merged_doc_text[:2500]}\n\n"
                    f"Question:\n{clean_q}"
                )
                ai_enhanced_resp = await self._call_gemini(ai_enhanced_prompt)
                ai_raw_answer = ai_enhanced_resp if (ai_enhanced_resp and len(ai_enhanced_resp.strip()) > 20) else file_raw_answer

        if chosen_source == "ai":
            # Pure AI Fallback Answer (Feature 6)
            ai_fallback_prompt = (
                "You are CampusX StudyGen, an expert educational and academic tutor. "
                "Provide a complete, detailed, and accurate explanation for the user's question.\n"
                "Rules:\n"
                f"{mark_guidelines}\n"
                f"{intent_guidelines}\n"
                "- For mathematical equations, format them clearly with standard LaTeX ($...$ for inline math, $$...$$ for block equations).\n"
                "- For chemical equations and formulas, format them clearly (e.g. CO₂ + H₂O → C₆H₁₂O₆ + O₂) with clear subscripts and arrows.\n"
                "- When presenting comparisons, properties, components, or tabular data, format them using clean Markdown tables.\n"
                "- Never ask the user to consult study materials or other sources. Always provide the actual, complete explanation directly.\n"
                "- Never output placeholder messages or error text.\n\n"
                f"Question:\n{clean_q}"
            )
            ai_raw_answer = await self._call_gemini(ai_fallback_prompt)

            if not ai_raw_answer:
                ai_raw_answer = (
                    "Unable to generate an AI answer due to a network or connection issue. "
                    "Please verify your connection and try again."
                )

        # Helper to strip accidental source suffixes and format long answers
        def _format_ans(ans: str) -> str:
            if not ans:
                return ""
            c = re.sub(
                r'[\r\n]+(?:\*{0,2}Source:?\*{0,2}\s*)?(?:[-*•]\s*)?(?:📄\s*|🤖\s*)?(?:File|Document|AI|Gemini)[\s_-]*Based[\s.]*$',
                '', ans, flags=re.IGNORECASE
            ).strip()
            return format_and_split_long_answer(c, max_chars=3800)

        final_file_answer = _format_ans(file_raw_answer) if chosen_source == "document" else None
        final_ai_answer = _format_ans(ai_raw_answer)

        final_response = {
            "answer": final_file_answer if chosen_source == "document" else final_ai_answer,
            "fileAnswer": final_file_answer,
            "aiAnswer": final_ai_answer,
            "source": chosen_source,
            "sourceDoc": selected_source_doc,
            "fileType": matched_file_type,
            "marks": detected_marks,
            "intent": detected_intent
        }

        self.response_cache[cache_key] = final_response
        return final_response


rag_service = RAGService()

